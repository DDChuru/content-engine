"""
Generate Ecowize Technical Architecture PowerPoint
e-wizer Digital Cleaning Operations Platform — Technical Companion Deck
For: Kurt Yon, Head of IT, Ecowize South Africa
From: Orlicron
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Brand Colors ──────────────────────────────────────────────────────
BLUE = RGBColor(0x06, 0x93, 0xE3)       # e-wizer primary
DARK = RGBColor(0x32, 0x37, 0x3C)       # Ecowize dark
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF0, 0xF4, 0xF8)
MID_GRAY = RGBColor(0x8B, 0x94, 0x9E)
GREEN = RGBColor(0x00, 0xD0, 0x84)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
DARK_BG = RGBColor(0x0D, 0x1B, 0x2A)
SECTION_BG = RGBColor(0x14, 0x23, 0x33)
PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
CYAN = RGBColor(0x06, 0xB6, 0xD4)

# ── Asset Paths ───────────────────────────────────────────────────────
ASSET_DIR = '/home/dachu/Documents/projects/angular/ncr_audit_app/src/assets/ecowize/images'
ECOWIZE_LOGO = '/tmp/ecowize-logo.png'
EWIZER_LOGO = f'{ASSET_DIR}/ewizer-logo.png'
HERO_BG = f'{ASSET_DIR}/ecowize-hero-v4.png'

# ── Presentation Setup ────────────────────────────────────────────────
TOTAL_SLIDES = 11

prs = Presentation()
prs.slide_width = Inches(13.333)   # Widescreen 16:9
prs.slide_height = Inches(7.5)
SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height


def add_bg(slide, color):
    """Fill slide background with solid color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color, alpha=None):
    """Add a colored rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if alpha is not None:
        from lxml import etree
        nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
        spPr = shape._element.find('.//a:spPr', nsmap)
        if spPr is not None:
            solidFill = spPr.find('.//a:solidFill', nsmap)
            if solidFill is not None:
                srgb = solidFill.find('a:srgbClr', nsmap)
                if srgb is not None:
                    alpha_el = etree.SubElement(srgb, '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha')
                    alpha_el.set('val', str(int(alpha * 1000)))
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name='Calibri', line_spacing=1.2):
    """Add a text box with formatted text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    if line_spacing != 1.0:
        p.line_spacing = Pt(font_size * line_spacing)
    return txBox


def add_rich_text_box(slide, left, top, width, height):
    """Add an empty text box for rich text building."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    return txBox, tf


def add_paragraph(tf, text, font_size=14, color=WHITE, bold=False,
                  alignment=PP_ALIGN.LEFT, font_name='Calibri',
                  space_after=6, space_before=0, bullet=False):
    """Add a paragraph to an existing text frame."""
    if len(tf.paragraphs) == 1 and tf.paragraphs[0].text == '':
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    if bullet:
        p.level = 0
        pPr = p._pPr
        if pPr is None:
            from lxml import etree
            pPr = etree.SubElement(p._p, '{http://schemas.openxmlformats.org/drawingml/2006/main}pPr')
        from lxml import etree
        buChar = etree.SubElement(pPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}buChar')
        buChar.set('char', '\u2022')
    return p


def add_footer(slide, text="Confidential  |  Orlicron  |  February 2026"):
    """Add subtle footer bar."""
    add_rect(slide, Inches(0), SLIDE_H - Inches(0.45), SLIDE_W, Inches(0.45), RGBColor(0x0A, 0x14, 0x20))
    add_text_box(slide, Inches(0.8), SLIDE_H - Inches(0.4), Inches(11), Inches(0.35),
                 text, font_size=8, color=MID_GRAY, alignment=PP_ALIGN.LEFT)


def add_slide_number(slide, num, total=TOTAL_SLIDES):
    add_text_box(slide, SLIDE_W - Inches(1.2), SLIDE_H - Inches(0.4), Inches(0.8), Inches(0.35),
                 f"{num}/{total}", font_size=8, color=MID_GRAY, alignment=PP_ALIGN.RIGHT)


def add_section_header(slide, section_label, title, bg_color=WHITE, text_color=DARK):
    """Standard section header: left accent bar + section label + title + divider."""
    if bg_color == WHITE:
        add_rect(slide, Inches(0), Inches(0), Inches(0.06), SLIDE_H, BLUE)

    add_text_box(slide, Inches(0.8), Inches(0.6), Inches(4), Inches(0.4),
                 section_label, font_size=11, color=BLUE, bold=True)

    add_text_box(slide, Inches(0.8), Inches(1.1), Inches(10), Inches(0.7),
                 title, font_size=30, color=text_color, bold=True)

    add_rect(slide, Inches(0.8), Inches(2.0), Inches(1.5), Inches(0.035), BLUE)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_bg(slide, DARK_BG)

# Hero background image (full bleed)
slide.shapes.add_picture(HERO_BG, Inches(0), Inches(0), SLIDE_W, SLIDE_H)

# Dark overlay for text readability
add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, RGBColor(0x0D, 0x1B, 0x2A), alpha=60)

# Ecowize logo (top left)
slide.shapes.add_picture(ECOWIZE_LOGO, Inches(0.8), Inches(0.6), Inches(2.4))

# e-wizer logo (top right)
slide.shapes.add_picture(EWIZER_LOGO, SLIDE_W - Inches(2.2), Inches(0.4), Inches(1.4), Inches(1.4))

# Title block
add_text_box(slide, Inches(0.8), Inches(2.6), Inches(10), Inches(0.6),
             "TECHNICAL OVERVIEW", font_size=14, color=BLUE, bold=True,
             font_name='Calibri')

add_text_box(slide, Inches(0.8), Inches(3.1), Inches(10), Inches(1.2),
             "Technical Architecture Overview", font_size=40, color=WHITE, bold=True,
             font_name='Calibri', line_spacing=1.1)

add_text_box(slide, Inches(0.8), Inches(4.4), Inches(8), Inches(0.8),
             "e-wizer Digital Cleaning Operations Platform",
             font_size=20, color=MID_GRAY, bold=False, font_name='Calibri')

# Divider line
add_rect(slide, Inches(0.8), Inches(5.5), Inches(2), Inches(0.04), BLUE)

# Prepared for line
add_text_box(slide, Inches(0.8), Inches(5.7), Inches(8), Inches(0.4),
             "Prepared for Kurt Yon, Head of IT  |  February 2026", font_size=12, color=MID_GRAY)

add_text_box(slide, Inches(0.8), Inches(6.2), Inches(6), Inches(0.4),
             "Confidential", font_size=10, color=RGBColor(0x4A, 0x55, 0x62))


# ══════════════════════════════════════════════════════════════════════
# SLIDE 2 — PLATFORM ARCHITECTURE
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, "ARCHITECTURE", "Modern, Scalable, Cloud-Native")

# 3 layer cards
layers = [
    {
        "title": "Web Application",
        "color": BLUE,
        "items": [
            "Angular 19 + TypeScript + Ionic 8 + Tailwind CSS",
            "Progressive Web App capable",
            "Responsive design for desktop, tablet, mobile browsers",
        ],
    },
    {
        "title": "Mobile Application",
        "color": GREEN,
        "items": [
            "React Native + Expo",
            "Native iOS & Android",
            "Offline-capable with sync",
            "Camera, GPS, QR scanning built-in",
        ],
    },
    {
        "title": "Backend & Database",
        "color": AMBER,
        "items": [
            "Firebase Functions v2 (Node.js 20)",
            "Cloud Firestore (NoSQL, real-time)",
            "Firebase Authentication",
            "Google Cloud Platform infrastructure",
        ],
    },
]

card_w = Inches(3.7)
card_h = Inches(4.0)
card_gap = Inches(0.35)
card_start_x = Inches(0.8)
card_y = Inches(2.6)

for i, layer in enumerate(layers):
    x = card_start_x + i * (card_w + card_gap)

    # Card background
    add_rect(slide, x, card_y, card_w, card_h, RGBColor(0xFA, 0xFB, 0xFC))

    # Color top bar
    add_rect(slide, x, card_y, card_w, Inches(0.06), layer["color"])

    # Title
    add_text_box(slide, x + Inches(0.3), card_y + Inches(0.3), card_w - Inches(0.6), Inches(0.4),
                 layer["title"], font_size=16, color=DARK, bold=True)

    # Divider inside card
    add_rect(slide, x + Inches(0.3), card_y + Inches(0.8), Inches(0.8), Inches(0.025), layer["color"])

    # Bullet items
    txBox, tf = add_rich_text_box(slide, x + Inches(0.3), card_y + Inches(1.1), card_w - Inches(0.6), Inches(2.6))
    for j, item in enumerate(layer["items"]):
        add_paragraph(tf, f"\u2022  {item}",
                      font_size=12, color=RGBColor(0x4A, 0x55, 0x62),
                      space_after=10)

add_footer(slide)
add_slide_number(slide, 2)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 3 — SECURITY & DATA PROTECTION
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, "SECURITY", "Enterprise-Grade Protection at Every Layer")

# 4 security pillar cards in 2x2 grid
security_pillars = [
    {
        "title": "Authentication",
        "color": BLUE,
        "items": [
            "Firebase Auth with email/password",
            "SSO-ready (Google, Microsoft AD)",
            "JWT token-based sessions",
            "MFA support available",
        ],
    },
    {
        "title": "Authorization",
        "color": GREEN,
        "items": [
            "Role-based access control (RBAC)",
            "6 role tiers: Tenant Admin \u2192 Admin \u2192 Site Admin \u2192 Supervisor \u2192 Operator",
            "Each role sees only what they need",
        ],
    },
    {
        "title": "Data Isolation",
        "color": AMBER,
        "items": [
            "Company-scoped collections (CSC pattern)",
            "Every data path: companies/{companyId}/{collection}",
            "Firestore security rules enforce boundaries at database level",
            "Zero cross-tenant data leakage by design",
        ],
    },
    {
        "title": "Mobile Device Protection",
        "color": PURPLE,
        "items": [
            "Minimal data stored on device \u2014 app streams from cloud",
            "Lost or stolen device = no sensitive data exposed",
            "Token-based auth: remote session revocation",
            "TLS/SSL in transit, AES-256 at rest (Google-managed)",
            "POPIA-aligned data handling",
        ],
    },
]

sec_w = Inches(5.8)
sec_h = Inches(2.1)
sec_gap_x = Inches(0.4)
sec_gap_y = Inches(0.3)
sec_start_x = Inches(0.8)
sec_start_y = Inches(2.5)

for i, pillar in enumerate(security_pillars):
    col = i % 2
    row = i // 2
    x = sec_start_x + col * (sec_w + sec_gap_x)
    y = sec_start_y + row * (sec_h + sec_gap_y)

    # Card background
    add_rect(slide, x, y, sec_w, sec_h, RGBColor(0xF8, 0xFA, 0xFC))

    # Left accent bar
    add_rect(slide, x, y, Inches(0.05), sec_h, pillar["color"])

    # Title
    add_text_box(slide, x + Inches(0.25), y + Inches(0.15), sec_w - Inches(0.5), Inches(0.35),
                 pillar["title"], font_size=14, color=DARK, bold=True)

    # Items
    txBox, tf = add_rich_text_box(slide, x + Inches(0.25), y + Inches(0.55), sec_w - Inches(0.5), Inches(1.4))
    for item in pillar["items"]:
        add_paragraph(tf, f"\u2022  {item}",
                      font_size=10, color=RGBColor(0x64, 0x6E, 0x78),
                      space_after=4)

add_footer(slide)
add_slide_number(slide, 3)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 4 — MULTI-TENANCY FOR ECOWIZE CLIENT GROUPS
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, "MULTI-TENANCY", "One Platform, Scaled Across Your Client Base")

# Intro paragraph
add_text_box(slide, Inches(0.8), Inches(2.3), Inches(11.5), Inches(0.5),
             "Each Ecowize client group gets its own isolated environment with tailored configuration, "
             "branding, and reporting \u2014 all managed from a single platform.",
             font_size=13, color=RGBColor(0x4A, 0x55, 0x62), line_spacing=1.4)

# Example client group cards — illustrative, not revealing competitors
client_groups = [
    ("Client Group A", "Bakery Operations", "Custom SCI library, bakery-specific\nmonitoring schemes, ATP testing", BLUE),
    ("Client Group B", "Fresh Produce", "Cold-chain compliance, temperature\nmonitoring, produce-specific hygiene", GREEN),
    ("Client Group C", "Dairy Processing", "CIP protocols, allergen management,\ndairy hygiene standards", AMBER),
    ("Client Group D", "Meat Processing", "HACCP alignment, red-meat specific\nNCR categories, kill-floor zones", PURPLE),
]

cg_w = Inches(2.7)
cg_h = Inches(2.3)
cg_gap = Inches(0.35)
cg_start_x = Inches(0.8)
cg_y = Inches(3.1)

for i, (name, industry, desc, color) in enumerate(client_groups):
    x = cg_start_x + i * (cg_w + cg_gap)

    # Card background
    add_rect(slide, x, cg_y, cg_w, cg_h, RGBColor(0xFA, 0xFB, 0xFC))

    # Color top bar
    add_rect(slide, x, cg_y, cg_w, Inches(0.06), color)

    # Client group name
    add_text_box(slide, x + Inches(0.2), cg_y + Inches(0.25), cg_w - Inches(0.4), Inches(0.35),
                 name, font_size=14, color=DARK, bold=True)

    # Industry
    add_text_box(slide, x + Inches(0.2), cg_y + Inches(0.6), cg_w - Inches(0.4), Inches(0.25),
                 industry, font_size=10, color=color, bold=True)

    # Config details
    add_text_box(slide, x + Inches(0.2), cg_y + Inches(1.0), cg_w - Inches(0.4), Inches(1.1),
                 desc, font_size=10, color=RGBColor(0x64, 0x6E, 0x78), line_spacing=1.4)

# Key benefit box at bottom
key_y = Inches(5.8)
add_rect(slide, Inches(0.8), key_y, Inches(11.7), Inches(0.9), RGBColor(0xF0, 0xF7, 0xFF))
add_rect(slide, Inches(0.8), key_y, Inches(0.05), Inches(0.9), BLUE)

kp_txBox, kp_tf = add_rich_text_box(slide, Inches(1.1), key_y + Inches(0.12), Inches(11.0), Inches(0.7))
add_paragraph(kp_tf, "Scale Without Complexity", font_size=12, color=DARK, bold=True, space_after=4)
add_paragraph(kp_tf,
    "Add new client groups with independent configurations, branding, and user bases \u2014 without "
    "duplicating infrastructure. Each group\u2019s data is completely isolated at the database level.",
    font_size=11, color=RGBColor(0x64, 0x6E, 0x78), space_after=0)

add_footer(slide)
add_slide_number(slide, 4)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 5 — MOBILE PLATFORM
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, "MOBILE", "Native Performance, Enterprise Distribution")

# Two columns
col_w = Inches(5.6)
col_gap = Inches(0.6)
left_x = Inches(0.8)
right_x = left_x + col_w + col_gap
col_y = Inches(2.6)
col_h = Inches(4.0)

# Left column — App Capabilities
add_rect(slide, left_x, col_y, col_w, col_h, RGBColor(0xFA, 0xFB, 0xFC))
add_rect(slide, left_x, col_y, col_w, Inches(0.06), BLUE)

add_text_box(slide, left_x + Inches(0.3), col_y + Inches(0.25), col_w - Inches(0.6), Inches(0.35),
             "App Capabilities", font_size=16, color=DARK, bold=True)

capabilities = [
    "Camera integration (photo evidence)",
    "GPS / geolocation (proof of presence)",
    "QR code scanning (zone identification)",
    "Offline mode with automatic sync",
    "Push notifications",
    "OTA updates (no app store wait)",
]

cap_txBox, cap_tf = add_rich_text_box(slide, left_x + Inches(0.3), col_y + Inches(0.85), col_w - Inches(0.6), Inches(3.0))
for cap in capabilities:
    add_paragraph(cap_tf, f"\u2022  {cap}",
                  font_size=12, color=RGBColor(0x4A, 0x55, 0x62),
                  space_after=10)

# Right column — Distribution & Device Security
add_rect(slide, right_x, col_y, col_w, col_h, RGBColor(0xFA, 0xFB, 0xFC))
add_rect(slide, right_x, col_y, col_w, Inches(0.06), GREEN)

add_text_box(slide, right_x + Inches(0.3), col_y + Inches(0.25), col_w - Inches(0.6), Inches(0.35),
             "Distribution & Device Security", font_size=16, color=DARK, bold=True)

distribution = [
    "iOS App Store + Google Play Store",
    "Enterprise distribution (MDM compatible)",
    "Expo EAS Build for CI/CD",
    "React Native 0.81 (latest stable)",
    "Minimal data on device \u2014 cloud-first design",
    "Stolen device = no sensitive data at risk",
]

dist_txBox, dist_tf = add_rich_text_box(slide, right_x + Inches(0.3), col_y + Inches(0.85), col_w - Inches(0.6), Inches(3.0))
for item in distribution:
    add_paragraph(dist_tf, f"\u2022  {item}",
                  font_size=12, color=RGBColor(0x4A, 0x55, 0x62),
                  space_after=10)

add_footer(slide)
add_slide_number(slide, 5)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 6 — WHY CLOUD
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)

add_text_box(slide, Inches(0.8), Inches(0.6), Inches(4), Inches(0.4),
             "INFRASTRUCTURE", font_size=11, color=BLUE, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.1), Inches(10), Inches(0.7),
             "Why Cloud \u2014 Built for Always-On Availability",
             font_size=30, color=WHITE, bold=True)

add_rect(slide, Inches(0.8), Inches(2.0), Inches(1.5), Inches(0.035), BLUE)

# Cloud advantage narrative
cloud_items = [
    ("Always Available", "In food safety, information must be accessible at all times. Cloud ensures your "
     "audit data, compliance records, and cleaning schedules are available 24/7 from any location \u2014 "
     "critical during real recalls, mock recalls, and unannounced audits.", BLUE),
    ("No On-Site Servers", "Eliminates the risk of hardware failures, local network outages, or power loss "
     "taking your compliance data offline. No server rooms, no backup tapes, no maintenance windows.", GREEN),
    ("Real-Time Sync", "Changes made on the factory floor sync to management dashboards in milliseconds. "
     "Every connected device sees the same current data \u2014 no stale spreadsheets or overnight batch uploads.", AMBER),
    ("Automatic Scaling", "Whether you run 5 sites or 50, the platform scales automatically. No capacity "
     "planning, no infrastructure procurement, no migration projects.", PURPLE),
    ("99.95% Uptime SLA", "Google Cloud Platform\u2019s enterprise SLA. Firebase powers apps like "
     "Duolingo, Alibaba, and The New York Times \u2014 proven at massive scale.", CYAN),
    ("Data Protection Built In", "Daily automatic backups with point-in-time recovery. AES-256 encryption "
     "at rest. Global CDN for fast load times. Region selection for data residency compliance.",
     RGBColor(0xE1, 0x1D, 0x48)),
]

metric_w = Inches(3.7)
metric_h = Inches(1.35)
metric_gap_x = Inches(0.3)
metric_gap_y = Inches(0.25)
metric_start_x = Inches(0.8)
metric_start_y = Inches(2.6)

for i, (title, desc, color) in enumerate(cloud_items):
    col = i % 3
    row = i // 3
    x = metric_start_x + col * (metric_w + metric_gap_x)
    y = metric_start_y + row * (metric_h + metric_gap_y)

    add_rect(slide, x, y, metric_w, metric_h, SECTION_BG)
    add_rect(slide, x, y, Inches(0.04), metric_h, color)

    add_text_box(slide, x + Inches(0.2), y + Inches(0.12), metric_w - Inches(0.3), Inches(0.3),
                 title, font_size=13, color=WHITE, bold=True)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.5), metric_w - Inches(0.3), Inches(0.75),
                 desc, font_size=9, color=MID_GRAY, line_spacing=1.35)

# Bottom note
note_y = Inches(5.85)
add_rect(slide, Inches(0.8), note_y, Inches(11.7), Inches(0.7), SECTION_BG)
add_rect(slide, Inches(0.8), note_y, Inches(0.04), Inches(0.7), BLUE)

add_text_box(slide, Inches(1.1), note_y + Inches(0.15), Inches(11.0), Inches(0.5),
             "Our experience across multiple food safety clients has consistently confirmed that "
             "cloud availability is non-negotiable in environments where audits and recalls demand instant access to records.",
             font_size=11, color=RGBColor(0xA0, 0xAE, 0xC0), line_spacing=1.35)

add_footer(slide, "Confidential  |  Orlicron  |  February 2026")
add_slide_number(slide, 6)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 7 — INTEGRATION & APIs
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, "INTEGRATION", "Open Architecture, Ready to Connect")

# Integration options as cards in 2x3 grid
integrations = [
    ("REST APIs", "Standard HTTP endpoints for data exchange with existing Ecowize systems and third-party platforms.", BLUE),
    ("Webhook Support", "Event-driven notifications for real-time system-to-system communication on key events.", GREEN),
    ("PDF Reports", "Automated PDF report generation and email distribution for compliance documentation.", AMBER),
    ("Data Export", "CSV/Excel data export for all modules. Compatible with any reporting or analytics tool.", PURPLE),
    ("Firebase Admin SDK", "Programmatic access for custom integrations, batch operations, and administrative tasks.", CYAN),
    ("Future: API Gateway", "ERP/HR system connectivity via managed API gateway. Scoped as business needs evolve.", RGBColor(0xE1, 0x1D, 0x48)),
]

int_w = Inches(3.7)
int_h = Inches(1.6)
int_gap_x = Inches(0.3)
int_gap_y = Inches(0.3)
int_start_x = Inches(0.8)
int_start_y = Inches(2.5)

for i, (title, desc, color) in enumerate(integrations):
    col = i % 3
    row = i // 3
    x = int_start_x + col * (int_w + int_gap_x)
    y = int_start_y + row * (int_h + int_gap_y)

    add_rect(slide, x, y, int_w, int_h, RGBColor(0xF8, 0xFA, 0xFC))
    add_rect(slide, x, y, Inches(0.05), int_h, color)

    add_text_box(slide, x + Inches(0.25), y + Inches(0.15), int_w - Inches(0.4), Inches(0.3),
                 title, font_size=13, color=DARK, bold=True)
    add_text_box(slide, x + Inches(0.25), y + Inches(0.55), int_w - Inches(0.4), Inches(0.9),
                 desc, font_size=11, color=RGBColor(0x64, 0x6E, 0x78), line_spacing=1.4)

# Bottom note
note_y = Inches(6.0)
add_rect(slide, Inches(0.8), note_y, Inches(11.7), Inches(0.6), RGBColor(0xF0, 0xF7, 0xFF))
add_rect(slide, Inches(0.8), note_y, Inches(0.05), Inches(0.6), BLUE)

add_text_box(slide, Inches(1.1), note_y + Inches(0.12), Inches(11.0), Inches(0.4),
             "The platform is built with standard open protocols. Integration with existing Ecowize systems can be scoped as needed.",
             font_size=11, color=DARK)

add_footer(slide)
add_slide_number(slide, 7)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 8 — IP OWNERSHIP MODEL
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, "OWNERSHIP", "Your Platform, Your Code, Your Data")

# Two columns
own_left_x = Inches(0.8)
own_right_x = Inches(7.0)
own_w = Inches(5.6)
own_y = Inches(2.5)
own_h = Inches(3.6)

# Left — What Ecowize Receives
add_rect(slide, own_left_x, own_y, own_w, own_h, RGBColor(0xF0, 0xF7, 0xFF))
add_rect(slide, own_left_x, own_y, own_w, Inches(0.55), BLUE)
add_text_box(slide, own_left_x + Inches(0.3), own_y + Inches(0.1), own_w - Inches(0.6), Inches(0.4),
             "What Ecowize Receives", font_size=14, color=WHITE, bold=True)

receives = [
    "Complete source code (web + mobile + backend)",
    "Firebase project ownership (transferable)",
    "Database with all data",
    "Domain and hosting configuration",
    "Documentation and deployment guides",
    "No licensing fees, no per-user charges, no recurring IP costs",
]

recv_txBox, recv_tf = add_rich_text_box(slide, own_left_x + Inches(0.3), own_y + Inches(0.75),
                                         own_w - Inches(0.6), Inches(2.7))
for item in receives:
    add_paragraph(recv_tf, f"\u2022  {item}",
                  font_size=11, color=DARK, space_after=8)

# Right — What This Means
add_rect(slide, own_right_x, own_y, own_w, own_h, RGBColor(0xF0, 0xFD, 0xF4))
add_rect(slide, own_right_x, own_y, own_w, Inches(0.55), GREEN)
add_text_box(slide, own_right_x + Inches(0.3), own_y + Inches(0.1), own_w - Inches(0.6), Inches(0.4),
             "What This Means", font_size=14, color=WHITE, bold=True)

means = [
    "Ecowize can engage any developer to maintain or extend",
    "No vendor lock-in \u2014 full technical independence",
    "The R15,000/month support is optional, engaged at Ecowize\u2019s discretion",
    "Source code in Git repository, full commit history",
]

means_txBox, means_tf = add_rich_text_box(slide, own_right_x + Inches(0.3), own_y + Inches(0.75),
                                           own_w - Inches(0.6), Inches(2.7))
for item in means:
    add_paragraph(means_tf, f"\u2022  {item}",
                  font_size=11, color=DARK, space_after=10)

# Bottom highlight
hl_y = Inches(6.3)
add_rect(slide, Inches(0.8), hl_y, Inches(11.7), Inches(0.6), RGBColor(0xF0, 0xF7, 0xFF))
add_rect(slide, Inches(0.8), hl_y, Inches(0.05), Inches(0.6), BLUE)

add_text_box(slide, Inches(1.1), hl_y + Inches(0.12), Inches(11.0), Inches(0.4),
             "This is a platform purchase, not a subscription. Ecowize owns 100% of the intellectual property.",
             font_size=12, color=DARK, bold=True)

add_footer(slide)
add_slide_number(slide, 8)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 9 — OUR TEAM
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, "YOUR TEAM", "Dedicated Expertise, End to End")

# Team members
team = [
    {
        "name": "Daniel Churu",
        "role": "Project Lead & Domain Expert",
        "desc": "Food safety technology specialist with deep understanding of SHEQ operations, "
                "audit compliance, and the regulatory landscape. Leads client engagement, "
                "requirements, and change management.",
        "color": BLUE,
    },
    {
        "name": "Dumisani Maramba",
        "role": "Senior Developer",
        "desc": "Full-stack development lead across web and mobile platforms. "
                "Responsible for core feature delivery, code quality, and technical decision-making.",
        "color": GREEN,
    },
    {
        "name": "Tanaka Musendo",
        "role": "Solutions Architect",
        "desc": "Platform architecture, database design, cloud infrastructure, and security. "
                "Ensures scalability, performance, and multi-tenant data integrity.",
        "color": AMBER,
    },
    {
        "name": "Tafadzwa Maramba",
        "role": "Developer",
        "desc": "Frontend and backend development, feature implementation, and testing. "
                "Hands-on delivery across Angular, React Native, and Firebase.",
        "color": PURPLE,
    },
    {
        "name": "Dary Mambowatema",
        "role": "Technical Support",
        "desc": "First-line support, site onboarding, user training, and issue resolution. "
                "The bridge between the technical team and your operational staff on the ground.",
        "color": CYAN,
    },
]

# Layout: 3 cards top row, 2 cards bottom row (centered)
tm_w = Inches(3.7)
tm_h = Inches(2.05)
tm_gap = Inches(0.35)
tm_start_x = Inches(0.8)
tm_start_y = Inches(2.5)

for i, member in enumerate(team):
    if i < 3:
        row = 0
        col = i
        x = tm_start_x + col * (tm_w + tm_gap)
    else:
        row = 1
        col = i - 3
        # Center 2 cards in bottom row
        bottom_total = 2 * tm_w + tm_gap
        bottom_start = (prs.slide_width - bottom_total) / 2
        x = bottom_start + col * (tm_w + tm_gap)
    y = tm_start_y + row * (tm_h + Inches(0.3))

    # Card background
    add_rect(slide, x, y, tm_w, tm_h, RGBColor(0xFA, 0xFB, 0xFC))

    # Color left accent bar
    add_rect(slide, x, y, Inches(0.05), tm_h, member["color"])

    # Name
    add_text_box(slide, x + Inches(0.25), y + Inches(0.15), tm_w - Inches(0.5), Inches(0.35),
                 member["name"], font_size=15, color=DARK, bold=True)

    # Role badge
    add_text_box(slide, x + Inches(0.25), y + Inches(0.5), tm_w - Inches(0.5), Inches(0.25),
                 member["role"], font_size=10, color=member["color"], bold=True)

    # Description
    add_text_box(slide, x + Inches(0.25), y + Inches(0.9), tm_w - Inches(0.5), Inches(1.0),
                 member["desc"], font_size=10, color=RGBColor(0x64, 0x6E, 0x78), line_spacing=1.35)

add_footer(slide)
add_slide_number(slide, 9)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 10 — TRACK RECORD & APPROACH
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, "DELIVERY", "Proven at Scale, Built Through Partnership")

# Left column — Track Record
tr_left_x = Inches(0.8)
tr_right_x = Inches(7.0)
tr_w = Inches(5.6)
tr_y = Inches(2.5)
tr_h = Inches(4.0)

add_rect(slide, tr_left_x, tr_y, tr_w, tr_h, RGBColor(0xF0, 0xF7, 0xFF))
add_rect(slide, tr_left_x, tr_y, tr_w, Inches(0.55), BLUE)
add_text_box(slide, tr_left_x + Inches(0.3), tr_y + Inches(0.1), tr_w - Inches(0.6), Inches(0.4),
             "Delivery Track Record", font_size=14, color=WHITE, bold=True)

track_items = [
    "Multi-site food safety platform deployed across active operations",
    "Proven in environments requiring unannounced audit readiness",
    "Successfully digitised cleaning schedules, NCR workflows, and compliance reporting",
    "Real-world recall response support \u2014 mock and live scenarios",
    "Platform handles hundreds of daily inspections across multiple sites",
    "Cloned and configured new tenants within days, not months",
]

tr_txBox, tr_tf = add_rich_text_box(slide, tr_left_x + Inches(0.3), tr_y + Inches(0.75),
                                     tr_w - Inches(0.6), Inches(3.1))
for item in track_items:
    add_paragraph(tr_tf, f"\u2022  {item}",
                  font_size=11, color=DARK, space_after=8)

# Right column — Change Management Approach
add_rect(slide, tr_right_x, tr_y, tr_w, tr_h, RGBColor(0xF0, 0xFD, 0xF4))
add_rect(slide, tr_right_x, tr_y, tr_w, Inches(0.55), GREEN)
add_text_box(slide, tr_right_x + Inches(0.3), tr_y + Inches(0.1), tr_w - Inches(0.6), Inches(0.4),
             "Our Mobilisation Approach", font_size=14, color=WHITE, bold=True)

change_items = [
    ("Strategic Alignment", "Start with SHEQ leadership to align configurations with your existing standards and audit requirements"),
    ("Graduated Rollout", "Pilot site first, then phased expansion \u2014 learn and adapt before scaling"),
    ("On-the-Ground Training", "Hands-on, site-level training with operational staff, not just management"),
    ("Champion Model", "Identify and empower site champions who drive adoption from within the team"),
    ("Continuous Feedback Loop", "Regular check-ins, usage analytics, and iterative improvement based on real operational data"),
]

cm_txBox, cm_tf = add_rich_text_box(slide, tr_right_x + Inches(0.3), tr_y + Inches(0.75),
                                     tr_w - Inches(0.6), Inches(3.1))
for title, desc in change_items:
    add_paragraph(cm_tf, title, font_size=11, color=DARK, bold=True, space_after=2)
    add_paragraph(cm_tf, desc, font_size=10, color=RGBColor(0x64, 0x6E, 0x78), space_after=10)

add_footer(slide)
add_slide_number(slide, 10)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 11 — CLOSING
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Full dark background
add_bg(slide, DARK_BG)

# Subtle hero bg
slide.shapes.add_picture(HERO_BG, Inches(0), Inches(0), SLIDE_W, SLIDE_H)
add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, RGBColor(0x0D, 0x1B, 0x2A), alpha=75)

# Ecowize logo
slide.shapes.add_picture(ECOWIZE_LOGO, Inches(4.5), Inches(1.0), Inches(4.2))

# e-wizer logo
slide.shapes.add_picture(EWIZER_LOGO, Inches(5.8), Inches(2.4), Inches(1.6), Inches(1.6))

# Divider
add_rect(slide, Inches(5.5), Inches(4.2), Inches(2.3), Inches(0.035), BLUE)

# Main message
add_text_box(slide, Inches(1.5), Inches(4.2), Inches(10.3), Inches(0.8),
             "Ready for a Technical Deep-Dive",
             font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(2), Inches(5.0), Inches(9.3), Inches(0.7),
             "We welcome the opportunity to walk through the platform\n"
             "with your IT team.",
             font_size=14, color=MID_GRAY, alignment=PP_ALIGN.CENTER, line_spacing=1.4)

# Contact details — two cards side by side
contact_y = Inches(5.7)
contact_w = Inches(4.0)
contact_h = Inches(0.9)
contact_gap = Inches(0.5)
contacts_total = 2 * contact_w + contact_gap
contacts_start_x = (SLIDE_W - contacts_total) / 2

# Daniel
dx = contacts_start_x
add_rect(slide, dx, contact_y, contact_w, contact_h, RGBColor(0x14, 0x23, 0x33))
add_rect(slide, dx, contact_y, contact_w, Inches(0.035), BLUE)
add_text_box(slide, dx + Inches(0.25), contact_y + Inches(0.1), contact_w - Inches(0.5), Inches(0.3),
             "Daniel Churu", font_size=13, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, dx + Inches(0.25), contact_y + Inches(0.42), contact_w - Inches(0.5), Inches(0.2),
             "079 797 2072  |  daniel@orlicron.com", font_size=10, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# Dumisani
dx2 = contacts_start_x + contact_w + contact_gap
add_rect(slide, dx2, contact_y, contact_w, contact_h, RGBColor(0x14, 0x23, 0x33))
add_rect(slide, dx2, contact_y, contact_w, Inches(0.035), GREEN)
add_text_box(slide, dx2 + Inches(0.25), contact_y + Inches(0.1), contact_w - Inches(0.5), Inches(0.3),
             "Dumisani Maramba", font_size=13, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, dx2 + Inches(0.25), contact_y + Inches(0.42), contact_w - Inches(0.5), Inches(0.2),
             "078 486 4667  |  dumi@orlicron.com", font_size=10, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# Footer line
add_text_box(slide, Inches(3), Inches(6.85), Inches(7.3), Inches(0.3),
             "Orlicron  |  February 2026  |  Confidential",
             font_size=9, color=RGBColor(0x4A, 0x55, 0x62), alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════════════
output_dir = '/home/dachu/Documents/projects/content-engine/packages/backend/output/ecowize'
os.makedirs(output_dir, exist_ok=True)
output_path = os.path.join(output_dir, 'ecowize-technical-overview-feb2026.pptx')
prs.save(output_path)
print(f"\n  Saved: {output_path}")
print(f"  Slides: {len(prs.slides)}")
print(f"  Format: 16:9 Widescreen\n")
