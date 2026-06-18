"""
Generate Ecowize CFO Proposal PowerPoint
e-wizer Digital Cleaning Operations Platform
From: Orlicron
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Brand Colors ──────────────────────────────────────────────────────
BLUE = RGBColor(0x06, 0x93, 0xE3)       # e-wizer primary
DARK = RGBColor(0x32, 0x37, 0x3C)       # Ecowize dark
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF0, 0xF4, 0xF8)
MID_GRAY = RGBColor(0x8B, 0x94, 0x9E)
GREEN = RGBColor(0x00, 0xD0, 0x84)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
DARK_BG = RGBColor(0x0D, 0x1B, 0x2A)
SECTION_BG = RGBColor(0x14, 0x23, 0x33)

# ── Asset Paths ───────────────────────────────────────────────────────
ASSET_DIR = '/home/dachu/Documents/projects/angular/ncr_audit_app/src/assets/ecowize/images'
ECOWIZE_LOGO = '/tmp/ecowize-logo.png'
EWIZER_LOGO = f'{ASSET_DIR}/ewizer-logo.png'
HERO_BG = f'{ASSET_DIR}/ecowize-hero-v4.png'
PILLARS_IMG = f'{ASSET_DIR}/ewizer-opex-pillars.png'
CLEANING_IMG = f'{ASSET_DIR}/ewizer-module-cleaning.png'
DASHBOARD_IMG = f'{ASSET_DIR}/ewizer-module-dashboard.png'

# ── Presentation Setup ────────────────────────────────────────────────
prs = Presentation()
prs.slide_width = Inches(13.333)   # Widescreen 16:9
prs.slide_height = Inches(7.5)
SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height


def add_bg(slide, color):
    """Fill slide background with solid color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color, alpha=None):
    """Add a colored rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if alpha is not None:
        from lxml import etree
        nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
        spPr = shape._element.find('.//a:spPr', nsmap)
        if spPr is not None:
            solidFill = spPr.find('.//a:solidFill', nsmap)
            if solidFill is not None:
                srgb = solidFill.find('a:srgbClr', nsmap)
                if srgb is not None:
                    alpha_el = etree.SubElement(srgb, '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha')
                    alpha_el.set('val', str(int(alpha * 1000)))
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name='Calibri', line_spacing=1.2):
    """Add a text box with formatted text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    if line_spacing != 1.0:
        p.line_spacing = Pt(font_size * line_spacing)
    return txBox


def add_rich_text_box(slide, left, top, width, height):
    """Add an empty text box for rich text building."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    return txBox, tf


def add_paragraph(tf, text, font_size=14, color=WHITE, bold=False,
                  alignment=PP_ALIGN.LEFT, font_name='Calibri',
                  space_after=6, space_before=0, bullet=False):
    """Add a paragraph to an existing text frame."""
    if len(tf.paragraphs) == 1 and tf.paragraphs[0].text == '':
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    if bullet:
        p.level = 0
        pPr = p._pPr
        if pPr is None:
            from lxml import etree
            pPr = etree.SubElement(p._p, '{http://schemas.openxmlformats.org/drawingml/2006/main}pPr')
        from lxml import etree
        buChar = etree.SubElement(pPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}buChar')
        buChar.set('char', '\u2022')
    return p


def add_footer(slide, text="Confidential  |  Orlicron  |  February 2026"):
    """Add subtle footer bar."""
    add_rect(slide, Inches(0), SLIDE_H - Inches(0.45), SLIDE_W, Inches(0.45), RGBColor(0x0A, 0x14, 0x20))
    add_text_box(slide, Inches(0.8), SLIDE_H - Inches(0.4), Inches(11), Inches(0.35),
                 text, font_size=8, color=MID_GRAY, alignment=PP_ALIGN.LEFT)


def add_slide_number(slide, num, total=9):
    add_text_box(slide, SLIDE_W - Inches(1.2), SLIDE_H - Inches(0.4), Inches(0.8), Inches(0.35),
                 f"{num}/{total}", font_size=8, color=MID_GRAY, alignment=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_bg(slide, DARK_BG)

# Hero background image (full bleed)
slide.shapes.add_picture(HERO_BG, Inches(0), Inches(0), SLIDE_W, SLIDE_H)

# Dark overlay for text readability
add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, RGBColor(0x0D, 0x1B, 0x2A), alpha=60)

# Ecowize logo (top left)
slide.shapes.add_picture(ECOWIZE_LOGO, Inches(0.8), Inches(0.6), Inches(2.4))

# e-wizer logo (top right)
slide.shapes.add_picture(EWIZER_LOGO, SLIDE_W - Inches(2.2), Inches(0.4), Inches(1.4), Inches(1.4))

# Title block
add_text_box(slide, Inches(0.8), Inches(2.6), Inches(10), Inches(0.6),
             "PROPOSAL", font_size=14, color=BLUE, bold=True,
             font_name='Calibri')

add_text_box(slide, Inches(0.8), Inches(3.1), Inches(10), Inches(1.2),
             "Digital Cleaning Operations Platform", font_size=40, color=WHITE, bold=True,
             font_name='Calibri', line_spacing=1.1)

add_text_box(slide, Inches(0.8), Inches(4.4), Inches(8), Inches(0.8),
             "Always-On Due Diligence for Ecowize South Africa",
             font_size=20, color=MID_GRAY, bold=False, font_name='Calibri')

# Divider line
add_rect(slide, Inches(0.8), Inches(5.5), Inches(2), Inches(0.04), BLUE)

# From line
add_text_box(slide, Inches(0.8), Inches(5.7), Inches(6), Inches(0.4),
             "Prepared by Orlicron  |  February 2026", font_size=12, color=MID_GRAY)

add_text_box(slide, Inches(0.8), Inches(6.2), Inches(6), Inches(0.4),
             "Confidential", font_size=10, color=RGBColor(0x4A, 0x55, 0x62))


# ══════════════════════════════════════════════════════════════════════
# SLIDE 2 — THE OPPORTUNITY
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

# Left accent bar
add_rect(slide, Inches(0), Inches(0), Inches(0.06), SLIDE_H, BLUE)

# Section label
add_text_box(slide, Inches(0.8), Inches(0.6), Inches(4), Inches(0.4),
             "THE OPPORTUNITY", font_size=11, color=BLUE, bold=True)

# Main heading
add_text_box(slide, Inches(0.8), Inches(1.1), Inches(7), Inches(1.0),
             "Transforming Cleaning Operations\ninto a Digital, Proactive System",
             font_size=30, color=DARK, bold=True, line_spacing=1.15)

# Divider
add_rect(slide, Inches(0.8), Inches(2.4), Inches(1.5), Inches(0.035), BLUE)

# Body text
body_txBox, body_tf = add_rich_text_box(slide, Inches(0.8), Inches(2.8), Inches(6.5), Inches(3.5))

add_paragraph(body_tf,
    "Ecowize South Africa delivers food safety and sanitation excellence across multiple sites "
    "and clients. The operational reality is clear: site teams execute cleaning, but Ecowize "
    "management holds the accountability.",
    font_size=14, color=DARK, space_after=14)

add_paragraph(body_tf,
    "e-wizer is a digital platform built specifically for Ecowize to provide always-on "
    "visibility into cleaning operations \u2014 giving every level of the business the information "
    "they need, when they need it.",
    font_size=14, color=DARK, space_after=14)

add_paragraph(body_tf,
    "This proposal focuses on the Cleaning Operation. The architecture is designed to "
    "extend into Pest Control, Hygiene, and any other operational area as the partnership grows.",
    font_size=14, color=RGBColor(0x64, 0x6E, 0x78), space_after=14)

# Right side — Dashboard image
slide.shapes.add_picture(DASHBOARD_IMG, Inches(8.0), Inches(1.0), Inches(4.8), Inches(2.7))

# Key message box at bottom
msg_rect = add_rect(slide, Inches(8.0), Inches(4.2), Inches(4.8), Inches(2.2), RGBColor(0xF0, 0xF7, 0xFF))
msg_txBox, msg_tf = add_rich_text_box(slide, Inches(8.3), Inches(4.4), Inches(4.2), Inches(2.0))
add_paragraph(msg_tf, "Key Principle", font_size=11, color=BLUE, bold=True, space_after=8)
add_paragraph(msg_tf,
    "The solution becomes Ecowize intellectual property. All source code is owned by Ecowize. "
    "e-wizer is your platform \u2014 a competitive advantage that belongs to you.",
    font_size=12, color=DARK, space_after=0)

add_footer(slide)
add_slide_number(slide, 2)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 3 — VALUE AT EVERY LEVEL
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_rect(slide, Inches(0), Inches(0), Inches(0.06), SLIDE_H, BLUE)

add_text_box(slide, Inches(0.8), Inches(0.6), Inches(4), Inches(0.4),
             "VALUE AT EVERY LEVEL", font_size=11, color=BLUE, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.1), Inches(10), Inches(0.7),
             "What e-wizer Means for Your Teams",
             font_size=30, color=DARK, bold=True)

add_rect(slide, Inches(0.8), Inches(2.0), Inches(1.5), Inches(0.035), BLUE)

# 5 value cards in a row
levels = [
    ("Site Teams", "Digital task lists, QR scanning, photo evidence. Know exactly what to clean, when, and confirm completion in real-time.", BLUE),
    ("Area Managers", "Live completion dashboards per zone. Exception alerts for missed tasks. Proof of work without site visits.", RGBColor(0x00, 0xD0, 0x84)),
    ("General Managers", "Multi-site compliance overview. Trend analysis across all operations. Faster decisions, fewer surprises.", RGBColor(0x8B, 0x5C, 0xF6)),
    ("SHEQ Team", "Audit-ready records at all times. NCR tracking with defined response times. Traceability from task to verification.", AMBER),
    ("Clients", "Confidence in cleaning compliance. Timestamped, geolocated proof of execution. Digital due diligence on demand.", RGBColor(0x06, 0xB6, 0xD4)),
]

card_w = Inches(2.25)
card_h = Inches(3.8)
start_x = Inches(0.6)
gap = Inches(0.18)
card_y = Inches(2.5)

for i, (title, desc, color) in enumerate(levels):
    x = start_x + i * (card_w + gap)

    # Card background
    card = add_rect(slide, x, card_y, card_w, card_h, RGBColor(0xFA, 0xFB, 0xFC))

    # Color top bar
    add_rect(slide, x, card_y, card_w, Inches(0.06), color)

    # Title
    add_text_box(slide, x + Inches(0.2), card_y + Inches(0.3), card_w - Inches(0.4), Inches(0.4),
                 title, font_size=14, color=DARK, bold=True)

    # Description
    add_text_box(slide, x + Inches(0.2), card_y + Inches(0.85), card_w - Inches(0.4), Inches(2.7),
                 desc, font_size=11, color=RGBColor(0x64, 0x6E, 0x78), line_spacing=1.4)

add_footer(slide)
add_slide_number(slide, 3)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 4 — SOLUTION SCOPE (Cleaning Focus)
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_rect(slide, Inches(0), Inches(0), Inches(0.06), SLIDE_H, BLUE)

add_text_box(slide, Inches(0.8), Inches(0.6), Inches(4), Inches(0.4),
             "SOLUTION SCOPE", font_size=11, color=BLUE, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.1), Inches(10), Inches(0.7),
             "e-wizer Cleaning Operations Platform",
             font_size=30, color=DARK, bold=True)

add_rect(slide, Inches(0.8), Inches(2.0), Inches(1.5), Inches(0.035), BLUE)

# 6 modules in 2x3 grid
modules = [
    ("Cleaning Verification", "QR-scanned zones with photo evidence. Exception-based reporting. Daily, weekly, and monthly cycles tracked independently."),
    ("NCR & Issue Tracking", "Four severity levels with defined response times. Seven-status workflow. Automatic escalation ensures nothing gets lost."),
    ("Surveys", "Configurable surveys for client satisfaction, internal quality checks, and operational feedback. Results feed directly into dashboards."),
    ("Training Management", "Staff competency tracking, training matrices, and compliance records per site. Ensure every team member is qualified."),
    ("Document & Record Control", "Centralised SOPs, work instructions, and cleaning specifications. Version-controlled with approval workflows."),
    ("Traceability", "Complete audit trail from task assignment to verification. Timestamped, geolocated, photo-verified. If it happened, it\u2019s recorded."),
]

cols = 3
rows = 2
mod_w = Inches(3.7)
mod_h = Inches(1.9)
mod_gap_x = Inches(0.3)
mod_gap_y = Inches(0.3)
mod_start_x = Inches(0.8)
mod_start_y = Inches(2.5)

for i, (title, desc) in enumerate(modules):
    col = i % cols
    row = i // cols
    x = mod_start_x + col * (mod_w + mod_gap_x)
    y = mod_start_y + row * (mod_h + mod_gap_y)

    # Card
    add_rect(slide, x, y, mod_w, mod_h, RGBColor(0xF8, 0xFA, 0xFC))
    add_rect(slide, x, y, Inches(0.05), mod_h, BLUE if i < 3 else GREEN)

    # Module title
    add_text_box(slide, x + Inches(0.25), y + Inches(0.2), mod_w - Inches(0.4), Inches(0.35),
                 title, font_size=13, color=DARK, bold=True)

    # Module description
    add_text_box(slide, x + Inches(0.25), y + Inches(0.65), mod_w - Inches(0.4), Inches(1.1),
                 desc, font_size=10, color=RGBColor(0x64, 0x6E, 0x78), line_spacing=1.4)

add_footer(slide)
add_slide_number(slide, 4)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 5 — INVESTMENT SUMMARY (The Numbers)
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_rect(slide, Inches(0), Inches(0), Inches(0.06), SLIDE_H, BLUE)

add_text_box(slide, Inches(0.8), Inches(0.6), Inches(4), Inches(0.4),
             "INVESTMENT SUMMARY", font_size=11, color=BLUE, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.1), Inches(10), Inches(0.7),
             "Clear, Transparent Pricing",
             font_size=30, color=DARK, bold=True)

add_rect(slide, Inches(0.8), Inches(2.0), Inches(1.5), Inches(0.035), BLUE)

# ── LEFT: Once-Off Investment ──
once_x = Inches(0.8)
once_y = Inches(2.5)
once_w = Inches(5.8)

# Once-off header
add_rect(slide, once_x, once_y, once_w, Inches(0.55), DARK)
add_text_box(slide, once_x + Inches(0.3), once_y + Inches(0.1), Inches(4), Inches(0.4),
             "Once-Off Investment", font_size=14, color=WHITE, bold=True)

# Table rows for once-off
once_items = [
    ("e-wizer Platform \u2014 Solution Build + Source Code Ownership", "R 485,000", True),
    ("QR Code Printer", "R 7,000", False),
    ("Cartridges", "Variable", False),
]

for j, (item, amount, highlight) in enumerate(once_items):
    row_y = once_y + Inches(0.55) + j * Inches(0.55)
    bg_color = RGBColor(0xF0, 0xF7, 0xFF) if highlight else RGBColor(0xF8, 0xFA, 0xFC)
    add_rect(slide, once_x, row_y, once_w, Inches(0.55), bg_color)

    # Item name
    fc = DARK if highlight else RGBColor(0x64, 0x6E, 0x78)
    fb = highlight
    add_text_box(slide, once_x + Inches(0.3), row_y + Inches(0.1), Inches(3.8), Inches(0.4),
                 item, font_size=11, color=fc, bold=fb)

    # Amount
    add_text_box(slide, once_x + Inches(4.0), row_y + Inches(0.1), Inches(1.5), Inches(0.4),
                 amount, font_size=13 if highlight else 11,
                 color=BLUE if highlight else DARK,
                 bold=highlight, alignment=PP_ALIGN.RIGHT)

# Note under once-off
note_y = once_y + Inches(0.55) + 3 * Inches(0.55) + Inches(0.15)
add_text_box(slide, once_x + Inches(0.3), note_y, Inches(5.2), Inches(0.6),
             "Cartridge cost varies based on site implementation rate.\n"
             "R485,000 includes complete solution build and full source code transfer to Ecowize.",
             font_size=9, color=MID_GRAY)

# ── RIGHT: Monthly Recurring ──
rec_x = Inches(7.2)
rec_y = Inches(2.5)
rec_w = Inches(5.4)

add_rect(slide, rec_x, rec_y, rec_w, Inches(0.55), RGBColor(0x14, 0x53, 0x2D))
add_text_box(slide, rec_x + Inches(0.3), rec_y + Inches(0.1), Inches(4), Inches(0.4),
             "Monthly Recurring", font_size=14, color=WHITE, bold=True)

recurring_items = [
    ("Database Infrastructure (Convex)", "R 5,000"),
    ("AI Services (Anthropic)", "R 6,000"),
    ("Support & Customisation", "R 15,000"),
]

for j, (item, amount) in enumerate(recurring_items):
    row_y = rec_y + Inches(0.55) + j * Inches(0.55)
    bg_color = RGBColor(0xF0, 0xFD, 0xF4) if j == 2 else RGBColor(0xF8, 0xFA, 0xFC)
    add_rect(slide, rec_x, row_y, rec_w, Inches(0.55), bg_color)

    add_text_box(slide, rec_x + Inches(0.3), row_y + Inches(0.1), Inches(3.2), Inches(0.4),
                 item, font_size=11, color=DARK)

    add_text_box(slide, rec_x + Inches(3.5), row_y + Inches(0.1), Inches(1.6), Inches(0.4),
                 amount, font_size=11, color=DARK, bold=True, alignment=PP_ALIGN.RIGHT)

# Total monthly
total_y = rec_y + Inches(0.55) + 3 * Inches(0.55)
add_rect(slide, rec_x, total_y, rec_w, Inches(0.55), BLUE)
add_text_box(slide, rec_x + Inches(0.3), total_y + Inches(0.1), Inches(3.2), Inches(0.4),
             "Total Monthly", font_size=12, color=WHITE, bold=True)
add_text_box(slide, rec_x + Inches(3.5), total_y + Inches(0.1), Inches(1.6), Inches(0.4),
             "R 26,000", font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.RIGHT)

# Note under recurring
rec_note_y = total_y + Inches(0.7)
add_text_box(slide, rec_x + Inches(0.3), rec_note_y, Inches(4.8), Inches(0.8),
             "Support includes customisations, issue resolution, site implementation assistance, "
             "and client-specific configurations. Engaged at Ecowize\u2019s discretion.",
             font_size=9, color=MID_GRAY)

# Bottom highlight box
highlight_y = Inches(5.8)
add_rect(slide, Inches(0.8), highlight_y, Inches(11.7), Inches(0.8), RGBColor(0xF0, 0xF7, 0xFF))
add_rect(slide, Inches(0.8), highlight_y, Inches(0.05), Inches(0.8), BLUE)

hl_txBox, hl_tf = add_rich_text_box(slide, Inches(1.1), highlight_y + Inches(0.15), Inches(11), Inches(0.6))
add_paragraph(hl_tf,
    "R485,000 is the total once-off investment for the complete e-wizer Cleaning Operations platform "
    "including full source code ownership by Ecowize. No licensing fees. No per-user charges. Your platform, your IP.",
    font_size=11, color=DARK, bold=False)

add_footer(slide)
add_slide_number(slide, 5)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 6 — IMPLEMENTATION APPROACH
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_rect(slide, Inches(0), Inches(0), Inches(0.06), SLIDE_H, BLUE)

add_text_box(slide, Inches(0.8), Inches(0.6), Inches(4), Inches(0.4),
             "IMPLEMENTATION", font_size=11, color=BLUE, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.1), Inches(10), Inches(0.7),
             "Fast to Deploy, Built to Last",
             font_size=30, color=DARK, bold=True)

add_rect(slide, Inches(0.8), Inches(2.0), Inches(1.5), Inches(0.035), BLUE)

# Three phases
phases = [
    {
        "num": "01",
        "title": "Strategic Alignment",
        "duration": "2 Sessions  |  Online",
        "desc": "Business strategic intent alignment for reporting and visualisation of operational excellence. "
                "We align on what the dashboards need to show and how the platform supports Ecowize\u2019s brand position.",
        "detail": "2 \u00d7 30min\u20131hr online sessions",
        "color": BLUE,
    },
    {
        "num": "02",
        "title": "SHEQ & System Configuration",
        "duration": "5 Days  |  Online Sessions",
        "desc": "Alignment of the Food Safety Management System to e-wizer. Initial report configuration, "
                "system baseline setup, cleaning schedules, zone mapping, and NCR workflows.",
        "detail": "5 working days with regular online collaboration",
        "color": GREEN,
    },
    {
        "num": "03",
        "title": "Site Rollout",
        "duration": "1 Day\u2019s Notice  |  Per Site",
        "desc": "Once SHEQ alignment is complete and the Master Cleaning Schedule is configured, "
                "site rollout is possible with as little as one day\u2019s notice. QR codes deployed, teams trained, system live.",
        "detail": "Rollout on MCS availability \u2014 rapid deployment",
        "color": AMBER,
    },
]

phase_w = Inches(3.7)
phase_h = Inches(3.6)
phase_gap = Inches(0.35)
phase_start_x = Inches(0.8)
phase_y = Inches(2.6)

for i, phase in enumerate(phases):
    x = phase_start_x + i * (phase_w + phase_gap)

    # Card background
    add_rect(slide, x, phase_y, phase_w, phase_h, RGBColor(0xFA, 0xFB, 0xFC))

    # Number circle (square with rounded look via rect)
    add_rect(slide, x + Inches(0.25), phase_y + Inches(0.3), Inches(0.55), Inches(0.55), phase["color"])
    add_text_box(slide, x + Inches(0.25), phase_y + Inches(0.33), Inches(0.55), Inches(0.5),
                 phase["num"], font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # Title
    add_text_box(slide, x + Inches(1.0), phase_y + Inches(0.3), Inches(2.5), Inches(0.4),
                 phase["title"], font_size=15, color=DARK, bold=True)

    # Duration badge
    add_text_box(slide, x + Inches(1.0), phase_y + Inches(0.7), Inches(2.5), Inches(0.3),
                 phase["duration"], font_size=10, color=phase["color"], bold=True)

    # Description
    add_text_box(slide, x + Inches(0.25), phase_y + Inches(1.3), phase_w - Inches(0.5), Inches(2.0),
                 phase["desc"], font_size=11, color=RGBColor(0x64, 0x6E, 0x78), line_spacing=1.4)

# Bottom note
add_text_box(slide, Inches(0.8), Inches(6.5), Inches(11), Inches(0.4),
             "We appreciate the urgency and are flexible to assist with delivery on any urgent requirements as we engage.",
             font_size=11, color=MID_GRAY, bold=False)

add_footer(slide)
add_slide_number(slide, 6)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 7 — KEY METRICS & ALWAYS-ON DUE DILIGENCE
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)

add_text_box(slide, Inches(0.8), Inches(0.6), Inches(4), Inches(0.4),
             "ALWAYS-ON DUE DILIGENCE", font_size=11, color=BLUE, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.1), Inches(10), Inches(0.7),
             "Measuring What Matters",
             font_size=30, color=WHITE, bold=True)

add_rect(slide, Inches(0.8), Inches(2.0), Inches(1.5), Inches(0.035), BLUE)

# Cleaning module image
slide.shapes.add_picture(CLEANING_IMG, Inches(8.5), Inches(0.5), Inches(4.3), Inches(2.4))

# Metrics cards
metrics = [
    ("Cleaning Completion Rate", "Real-time percentage of tasks completed per site, per shift. No more guessing."),
    ("NCR Response Time", "Track time from issue identification to resolution. Ensure critical NCRs are actioned within defined SLAs."),
    ("Compliance Score", "Site-level and area-level compliance percentages. Benchmark across the business."),
    ("Trend Analysis", "Week-over-week and month-over-month trends. See whether sites are improving or declining."),
    ("Proof of Presence", "Timestamped, geolocated, photo-verified evidence of cleaning execution. Digital proof, always available."),
    ("Audit Readiness", "Food safety management records accessible from any device, anywhere, at any time. Always audit-ready."),
]

metric_w = Inches(3.7)
metric_h = Inches(1.2)
metric_gap_x = Inches(0.3)
metric_gap_y = Inches(0.2)
metric_start_x = Inches(0.8)
metric_start_y = Inches(2.6)

for i, (title, desc) in enumerate(metrics):
    col = i % 3
    row = i // 3
    x = metric_start_x + col * (metric_w + metric_gap_x)
    y = metric_start_y + row * (metric_h + metric_gap_y)

    add_rect(slide, x, y, metric_w, metric_h, RGBColor(0x14, 0x23, 0x33))
    add_rect(slide, x, y, Inches(0.04), metric_h, BLUE if col == 0 else (GREEN if col == 1 else AMBER))

    add_text_box(slide, x + Inches(0.2), y + Inches(0.12), metric_w - Inches(0.3), Inches(0.3),
                 title, font_size=12, color=WHITE, bold=True)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.45), metric_w - Inches(0.3), Inches(0.7),
                 desc, font_size=9, color=MID_GRAY, line_spacing=1.3)

# Bottom quote
add_rect(slide, Inches(0.8), Inches(5.5), Inches(11.7), Inches(0.8), RGBColor(0x14, 0x23, 0x33))
add_text_box(slide, Inches(1.1), Inches(5.6), Inches(11), Inches(0.6),
             "\"An architecture that provides always-on access to Food Safety Management Records around cleaning \u2014 "
             "giving visibility across the entire business.\"",
             font_size=13, color=RGBColor(0xA0, 0xAE, 0xC0), bold=False)

add_footer(slide, "Confidential  |  Orlicron  |  February 2026")
add_slide_number(slide, 7)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 8 — BEYOND CLEANING (Architecture for Growth)
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_rect(slide, Inches(0), Inches(0), Inches(0.06), SLIDE_H, BLUE)

add_text_box(slide, Inches(0.8), Inches(0.6), Inches(4), Inches(0.4),
             "BUILT TO GROW", font_size=11, color=BLUE, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.1), Inches(10), Inches(0.7),
             "An Architecture That Extends",
             font_size=30, color=DARK, bold=True)

add_rect(slide, Inches(0.8), Inches(2.0), Inches(1.5), Inches(0.035), BLUE)

# Pillars image
slide.shapes.add_picture(PILLARS_IMG, Inches(7.5), Inches(1.8), Inches(5.2), Inches(2.9))

# Left text
body_txBox, body_tf = add_rich_text_box(slide, Inches(0.8), Inches(2.5), Inches(6.2), Inches(3.5))

add_paragraph(body_tf,
    "This phase delivers the Cleaning Operation.",
    font_size=15, color=DARK, bold=True, space_after=12)

add_paragraph(body_tf,
    "The e-wizer architecture is built on a foundation that allows extension into any operational area "
    "Ecowize chooses to digitise:",
    font_size=13, color=RGBColor(0x64, 0x6E, 0x78), space_after=16)

areas = [
    "Pest Control Operations",
    "Hygiene Services",
    "Waste Management",
    "Client-Specific Compliance Modules",
]
for area in areas:
    add_paragraph(body_tf, f"  \u2022  {area}",
                  font_size=13, color=DARK, space_after=8)

add_paragraph(body_tf, "", font_size=6, space_after=8)

add_paragraph(body_tf,
    "We are happy to engage with Ecowize on these areas as the partnership develops. "
    "The same platform, the same ownership model, the same always-on due diligence.",
    font_size=12, color=MID_GRAY, space_after=0)

# Five pillars labels at bottom
pillar_y = Inches(5.6)
pillar_labels = [
    ("Automation", BLUE),
    ("Observability", GREEN),
    ("Incident Management", AMBER),
    ("Continuous Improvement", RGBColor(0x8B, 0x5C, 0xF6)),
    ("Ownership", RGBColor(0x06, 0xB6, 0xD4)),
]
pillar_w = Inches(2.2)
for i, (label, color) in enumerate(pillar_labels):
    x = Inches(0.8) + i * (pillar_w + Inches(0.15))
    add_rect(slide, x, pillar_y, pillar_w, Inches(0.5), RGBColor(0xF8, 0xFA, 0xFC))
    add_rect(slide, x, pillar_y, pillar_w, Inches(0.04), color)
    add_text_box(slide, x + Inches(0.15), pillar_y + Inches(0.1), pillar_w - Inches(0.3), Inches(0.35),
                 label, font_size=10, color=DARK, bold=True, alignment=PP_ALIGN.CENTER)

add_footer(slide)
add_slide_number(slide, 8)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 9 — PARTNERSHIP (Closing)
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Full dark background
add_bg(slide, DARK_BG)

# Subtle hero bg
slide.shapes.add_picture(HERO_BG, Inches(0), Inches(0), SLIDE_W, SLIDE_H)
add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, RGBColor(0x0D, 0x1B, 0x2A), alpha=75)

# Ecowize logo
slide.shapes.add_picture(ECOWIZE_LOGO, Inches(4.5), Inches(1.0), Inches(4.2))

# e-wizer logo
slide.shapes.add_picture(EWIZER_LOGO, Inches(5.8), Inches(2.4), Inches(1.6), Inches(1.6))

# Divider
add_rect(slide, Inches(5.5), Inches(4.2), Inches(2.3), Inches(0.035), BLUE)

# Main message
add_text_box(slide, Inches(1.5), Inches(4.2), Inches(10.3), Inches(0.8),
             "Let\u2019s Build This Together",
             font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(2), Inches(5.0), Inches(9.3), Inches(0.7),
             "We are excited at the opportunity to partner with Ecowize and together deliver\n"
             "at an excellent level that matches your reputable brand position as a market leader\n"
             "in food safety, hygiene, and sanitation solutions.",
             font_size=13, color=MID_GRAY, alignment=PP_ALIGN.CENTER, line_spacing=1.35)

# Contact details — two cards side by side
contact_y = Inches(5.9)
contact_w = Inches(4.0)
contact_h = Inches(0.9)
contact_gap = Inches(0.5)
contacts_total = 2 * contact_w + contact_gap
contacts_start_x = (SLIDE_W - contacts_total) / 2

# Daniel
dx = contacts_start_x
add_rect(slide, dx, contact_y, contact_w, contact_h, RGBColor(0x14, 0x23, 0x33))
add_rect(slide, dx, contact_y, contact_w, Inches(0.035), BLUE)
add_text_box(slide, dx + Inches(0.25), contact_y + Inches(0.1), contact_w - Inches(0.5), Inches(0.3),
             "Daniel Churu", font_size=13, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, dx + Inches(0.25), contact_y + Inches(0.42), contact_w - Inches(0.5), Inches(0.2),
             "079 797 2072  |  daniel@orlicron.com", font_size=10, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# Dumisani
dx2 = contacts_start_x + contact_w + contact_gap
add_rect(slide, dx2, contact_y, contact_w, contact_h, RGBColor(0x14, 0x23, 0x33))
add_rect(slide, dx2, contact_y, contact_w, Inches(0.035), GREEN)
add_text_box(slide, dx2 + Inches(0.25), contact_y + Inches(0.1), contact_w - Inches(0.5), Inches(0.3),
             "Dumisani Maramba", font_size=13, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, dx2 + Inches(0.25), contact_y + Inches(0.42), contact_w - Inches(0.5), Inches(0.2),
             "078 486 4667  |  dumisani@orlicron.com", font_size=10, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# Footer line
add_text_box(slide, Inches(3), Inches(6.95), Inches(7.3), Inches(0.3),
             "Orlicron  |  February 2026  |  Confidential",
             font_size=9, color=RGBColor(0x4A, 0x55, 0x62), alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════════════
output_dir = '/home/dachu/Documents/projects/content-engine/packages/backend/output/ecowize'
os.makedirs(output_dir, exist_ok=True)
output_path = os.path.join(output_dir, 'ecowize-ewizer-proposal-feb2026.pptx')
prs.save(output_path)
print(f"\n  Saved: {output_path}")
print(f"  Slides: {len(prs.slides)}")
print(f"  Format: 16:9 Widescreen\n")
