#!/usr/bin/env python3
"""
PowerPoint Generator for Strategy Consultant Agent
Generates professional slide decks from strategic analysis data
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import json
import sys
from datetime import datetime


def set_text_format(text_frame, font_size=14, bold=False, color=None):
    """Apply consistent formatting to text"""
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.name = 'Calibri'
            if color:
                run.font.color.rgb = color


def add_title_slide(prs, data):
    """Create professional title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Background color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(31, 78, 121)  # Professional blue

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(3), Inches(14), Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_frame.text = data['title']
    set_text_format(title_frame, font_size=44, bold=True, color=RGBColor(255, 255, 255))
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(
        Inches(1), Inches(4.8), Inches(14), Inches(1)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_text = f"{data.get('company', '')} | {data.get('period', '')}"
    subtitle_frame.text = subtitle_text
    set_text_format(subtitle_frame, font_size=24, color=RGBColor(200, 200, 200))
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Date
    date_box = slide.shapes.add_textbox(
        Inches(1), Inches(8), Inches(14), Inches(0.5)
    )
    date_frame = date_box.text_frame
    date_frame.text = datetime.now().strftime("%B %d, %Y")
    set_text_format(date_frame, font_size=14, color=RGBColor(150, 150, 150))
    date_frame.paragraphs[0].alignment = PP_ALIGN.CENTER


def add_executive_summary_slide(prs, data):
    """Executive summary with key insights"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    title = slide.shapes.title
    title.text = "Executive Summary"

    body_shape = slide.placeholders[1]
    text_frame = body_shape.text_frame
    text_frame.clear()

    insights = data.get('insights', {})

    # Revenue insight
    if 'revenue' in insights:
        p = text_frame.add_paragraph()
        p.text = f"Revenue: ${insights['revenue'].get('current', 0):,}"
        p.level = 0
        if insights['revenue'].get('trend') == 'up':
            p.text += f" ↑ {insights['revenue'].get('change', 0)}%"

    # Profit margin insight
    if 'profitMargin' in insights:
        p = text_frame.add_paragraph()
        pm = insights['profitMargin']
        p.text = f"Profit Margin: {pm.get('current', 0)}% (Industry: {pm.get('benchmark', 0)}%)"
        p.level = 0

    # Cash flow insight
    if 'cashFlow' in insights:
        p = text_frame.add_paragraph()
        cf = insights['cashFlow']
        p.text = f"Cash Flow: ${cf.get('current', 0):,} ({cf.get('risk', 'medium')} risk)"
        p.level = 0


def add_financial_overview_slide(prs, data):
    """Financial overview with bar chart"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title only
    title = slide.shapes.title
    title.text = "Financial Performance"

    # Create chart data
    chart_data = CategoryChartData()

    quarters = data.get('quarters', ['Q1', 'Q2', 'Q3', 'Q4'])
    revenue_by_quarter = data.get('revenue_by_quarter', [400000, 435000, 465000, 500000])
    profit_by_quarter = data.get('profit_by_quarter', [75000, 85000, 95000, 110000])

    chart_data.categories = quarters
    chart_data.add_series('Revenue ($)', revenue_by_quarter)
    chart_data.add_series('Profit ($)', profit_by_quarter)

    # Add chart
    x, y, cx, cy = Inches(1.5), Inches(2), Inches(12), Inches(5)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    ).chart

    # Format chart
    chart.has_legend = True
    chart.legend.position = 2  # Bottom


def add_recommendations_slide(prs, data):
    """Strategic recommendations prioritized by impact"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    title = slide.shapes.title
    title.text = "Strategic Recommendations"

    body_shape = slide.placeholders[1]
    text_frame = body_shape.text_frame
    text_frame.clear()

    recommendations = data.get('recommendations', [])

    for i, rec in enumerate(recommendations[:5], 1):  # Top 5 recommendations
        p = text_frame.add_paragraph()
        p.text = f"{rec.get('category', 'Strategy').title()}"
        p.level = 0
        p.font.bold = True

        # Description
        p2 = text_frame.add_paragraph()
        p2.text = rec.get('description', '')
        p2.level = 1

        # Impact
        p3 = text_frame.add_paragraph()
        p3.text = f"Impact: {rec.get('impact', 'TBD')} | Effort: {rec.get('effort', 'medium')}"
        p3.level = 2
        p3.font.size = Pt(12)
        p3.font.italic = True


def add_customer_analysis_slide(prs, data):
    """Top customers with pie chart"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title only
    title = slide.shapes.title
    title.text = "Customer Revenue Distribution"

    insights = data.get('insights', {})
    top_customers = insights.get('topCustomers', [
        {'name': 'Enterprise Client A', 'revenue': 180000},
        {'name': 'Mid-Market Client B', 'revenue': 95000},
        {'name': 'Small Business C', 'revenue': 68000},
        {'name': 'Others', 'revenue': 157000}
    ])

    # Create pie chart data
    chart_data = CategoryChartData()
    categories = [c['name'] for c in top_customers]
    values = [c['revenue'] for c in top_customers]

    chart_data.categories = categories
    chart_data.add_series('Revenue', values)

    # Add chart
    x, y, cx, cy = Inches(2), Inches(2), Inches(12), Inches(5)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
    ).chart

    chart.has_legend = True
    chart.legend.position = 2  # Bottom


def add_next_steps_slide(prs, data):
    """Action items and timeline"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    title = slide.shapes.title
    title.text = "Next Steps"

    body_shape = slide.placeholders[1]
    text_frame = body_shape.text_frame
    text_frame.clear()

    next_steps = data.get('next_steps', [
        "Review and approve strategic recommendations",
        "Schedule implementation planning sessions",
        "Assign owners to priority initiatives",
        "Set up monthly performance review meetings",
        "Begin vendor consolidation analysis"
    ])

    for step in next_steps:
        p = text_frame.add_paragraph()
        p.text = step
        p.level = 0


def create_strategic_analysis(data, output_path):
    """Generate complete strategic analysis presentation"""
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    # Add all slides
    add_title_slide(prs, data)
    add_executive_summary_slide(prs, data)
    add_financial_overview_slide(prs, data)
    add_customer_analysis_slide(prs, data)
    add_recommendations_slide(prs, data)
    add_next_steps_slide(prs, data)

    # Save presentation
    prs.save(output_path)

    return {
        'success': True,
        'path': output_path,
        'slides': len(prs.slides),
        'template': 'strategic-analysis'
    }


def create_rfq_proposal(data, output_path):
    """Generate RFQ/Proposal presentation"""
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    # Title slide
    add_title_slide(prs, {
        'title': data.get('title', 'Business Proposal'),
        'company': data.get('client', ''),
        'period': data.get('date', datetime.now().strftime("%B %Y"))
    })

    # Executive Summary
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Executive Summary"

    body_shape = slide.placeholders[1]
    text_frame = body_shape.text_frame
    text_frame.clear()

    opportunity = data.get('opportunity', {})
    p = text_frame.add_paragraph()
    p.text = f"Service: {opportunity.get('service', 'Consulting Services')}"

    p = text_frame.add_paragraph()
    p.text = f"Timeline: {opportunity.get('timeline', '3-6 months')}"

    if opportunity.get('budget'):
        p = text_frame.add_paragraph()
        p.text = f"Investment: ${opportunity['budget']:,}"

    # Scope of Work
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Scope of Work"

    body_shape = slide.placeholders[1]
    text_frame = body_shape.text_frame
    text_frame.clear()

    deliverables = data.get('deliverables', [
        'Initial assessment and analysis',
        'Strategic recommendations',
        'Implementation roadmap',
        'Ongoing support and monitoring'
    ])

    for item in deliverables:
        p = text_frame.add_paragraph()
        p.text = item
        p.level = 0

    # Pricing (if included)
    if data.get('includePricing'):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = "Investment"

        body_shape = slide.placeholders[1]
        text_frame = body_shape.text_frame
        text_frame.clear()

        pricing = data.get('pricing', {})
        p = text_frame.add_paragraph()
        p.text = f"Professional Services: ${pricing.get('services', 35000):,}"

        p = text_frame.add_paragraph()
        p.text = f"Implementation Support: ${pricing.get('support', 10000):,}"

        p = text_frame.add_paragraph()
        p.text = f"Training: ${pricing.get('training', 5000):,}"

        p = text_frame.add_paragraph()
        p.text = ""

        p = text_frame.add_paragraph()
        total = sum([pricing.get('services', 35000), pricing.get('support', 10000), pricing.get('training', 5000)])
        p.text = f"Total Investment: ${total:,}"
        p.font.bold = True
        p.font.size = Pt(18)

    # Next Steps
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Next Steps"

    body_shape = slide.placeholders[1]
    text_frame = body_shape.text_frame
    text_frame.clear()

    p = text_frame.add_paragraph()
    p.text = "Review proposal"

    p = text_frame.add_paragraph()
    p.text = "Schedule kickoff meeting"

    p = text_frame.add_paragraph()
    p.text = "Sign agreement"

    p = text_frame.add_paragraph()
    p.text = "Begin engagement"

    # Save presentation
    prs.save(output_path)

    return {
        'success': True,
        'path': output_path,
        'slides': len(prs.slides),
        'template': 'rfq-proposal'
    }


def main():
    """Main entry point"""
    if len(sys.argv) < 3:
        print(json.dumps({
            'success': False,
            'error': 'Usage: generate-powerpoint.py <data_json> <output_path>'
        }))
        sys.exit(1)

    try:
        # Parse input
        data = json.loads(sys.argv[1])
        output_path = sys.argv[2]

        # Determine template type
        template = data.get('template', 'strategic-analysis')

        # Generate presentation
        if template == 'strategic-analysis':
            result = create_strategic_analysis(data, output_path)
        elif template == 'rfq-proposal':
            result = create_rfq_proposal(data, output_path)
        else:
            result = create_strategic_analysis(data, output_path)

        # Output result
        print(json.dumps(result))

    except Exception as e:
        print(json.dumps({
            'success': False,
            'error': str(e)
        }))
        sys.exit(1)


if __name__ == "__main__":
    main()
